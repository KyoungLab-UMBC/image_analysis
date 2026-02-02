import numpy as np
import matplotlib.pyplot as plt
import tifffile
from skimage import measure
from pptx import Presentation
from pptx.util import Inches
import io
import warnings
import os
from scipy import signal
from numba import jit
import concurrent.futures
from pptx.dml.color import RGBColor
import trend_analysis as trend_tools

# Set matplotlib backend to Agg for thread safety (no GUI)
plt.switch_backend('Agg')

# Suppress warnings
warnings.filterwarnings("ignore", category=RuntimeWarning)
warnings.filterwarnings("ignore", category=UserWarning)

# --- ACCELERATED WORKER FUNCTIONS (NUMBA) ---

@jit(nopython=True)
def calculate_spherical_profile(crop_r, crop_g, crop_mask, obj_id):
    """
    Numba-accelerated function to calculate radial profile in 3D.
    Iterates over the 21x21x21 cube and bins pixels by spherical radius.
    """
    d, h, w = crop_r.shape
    cz, cy, cx = d // 2, h // 2, w // 2
    
    # Radius 0 to 10 (11 bins)
    max_r = 11
    sums_r = np.zeros(max_r)
    sums_g = np.zeros(max_r)
    counts = np.zeros(max_r)
    
    for z in range(d):
        for y in range(h):
            for x in range(w):
                # Mask Logic: Include if background (0) or current object (obj_id)
                # Exclude if it belongs to a neighbor object
                if crop_mask[z, y, x] != 0 and crop_mask[z, y, x] != obj_id:
                    continue
                
                # 3D Distance
                dist = np.sqrt((z - cz)**2 + (y - cy)**2 + (x - cx)**2)
                
                # Floor to get bin index (0.5 -> 0, 1.9 -> 1)
                r_idx = int(dist)
                
                if r_idx < max_r:
                    sums_r[r_idx] += crop_r[z, y, x]
                    sums_g[r_idx] += crop_g[z, y, x]
                    counts[r_idx] += 1
                    
    means_r = np.zeros(max_r)
    means_g = np.zeros(max_r)
    
    for i in range(max_r):
        if counts[i] > 0:
            means_r[i] = sums_r[i] / counts[i]
            means_g[i] = sums_g[i] / counts[i]
            
    return means_r, means_g


def create_plot_image(red_data, green_data, x_axis, obj_id):
    fig, ax1 = plt.subplots(figsize=(6, 4))
    
    ax1.set_xlabel('Radius (pixel)')
    ax1.set_xlim(0, 10)
    ax1.set_xticks(range(0, 11, 1))
    
    # Left Y: Red
    ax1.set_ylabel('Norm. Intensity (Red)', color='red')
    # Use standard scaling or dynamic based on data
    ax1.plot(x_axis, red_data, color='red', linewidth=1.5)
    ax1.tick_params(axis='y', labelcolor='red')

    # Right Y: Green
    ax2 = ax1.twinx() 
    ax2.set_ylabel('Norm. Intensity (Green)', color='green')
    ax2.plot(x_axis, green_data, color='green', linewidth=1.5)
    ax2.tick_params(axis='y', labelcolor='green')

    if obj_id != "AVERAGE":
        plt.title(f"Object: {obj_id}")
    else:
        plt.title("Group Average (Normalized)")
    
    fig.tight_layout()
    img_stream = io.BytesIO()
    plt.savefig(img_stream, format='png', dpi=100)
    plt.close(fig)
    img_stream.seek(0)
    return img_stream

def create_dimetric_thumbnail(crop_r, crop_g):
    """
    Creates a 3D Volumetric Thumbnail.
    - Black Background
    - Alpha Scales with Intensity (Max 0.5) to prevent occlusion
    - Additive Blending (Red + Green = Yellow)
    - Custom Contrast (Median to Peak)
    """
    fig = plt.figure(figsize=(4, 4), facecolor='black')
    ax = fig.add_subplot(111, projection='3d')
    ax.set_facecolor('black')
    
    ax.grid(False)
    ax.set_axis_off()
    ax.xaxis.pane.fill = False
    ax.yaxis.pane.fill = False
    ax.zaxis.pane.fill = False

    s_start, s_end = 4, 17
    sub_r = crop_r[s_start:s_end, s_start:s_end, s_start:s_end]
    sub_g = crop_g[s_start:s_end, s_start:s_end, s_start:s_end]
    
    cz, cy, cx = 6, 6, 6 

    # --- NORMALIZE FUNCTION ---
    def get_normalized_channel(arr, is_red):
        v_min = np.median(arr)
        
        if is_red:
            center_val = arr[cz, cy, cx]
            v_max = center_val if center_val > v_min else np.max(arr)
            # Boost Red Exposure slightly to make mid-tones visible
            boost = 3.0 
        else:
            v_max = np.max(arr)
            boost = 1.0

        if v_max <= v_min: v_max = v_min + 1.0
            
        # Normalize
        arr_norm = (arr - v_min) / (v_max - v_min)
        arr_norm = arr_norm * boost
        return np.clip(arr_norm, 0, 1)

    norm_r = get_normalized_channel(sub_r, is_red=True)
    norm_g = get_normalized_channel(sub_g, is_red=False)

    # --- COMBINE & PLOT ---
    # 1. Calculate Combined Intensity (Max of R or G at each point)
    # This determines how "important" the pixel is.
    max_intensity = np.maximum(norm_r, norm_g)

    # 2. Threshold
    # We remove pixels that are too dark to matter (cleaner view)
    mask = max_intensity > 0.1 
    
    if np.any(mask):
        z, y, x = np.where(mask)
        
        colors = np.zeros((len(z), 4))
        colors[:, 0] = norm_r[mask] # Red
        colors[:, 1] = norm_g[mask] # Green
        colors[:, 2] = 0.0          # Blue
        
        # --- DYNAMIC ALPHA CORRECTION ---
        # Instead of fixed 0.5, we scale alpha by intensity.
        # Brightest pixel (1.0) -> 0.5 Alpha
        # Dim pixel (0.2)       -> 0.1 Alpha (Transparent)
        # This allows you to "see through" the dark parts.
        intensities = max_intensity[mask]
        colors[:, 3] = np.clip(intensities * 0.5, 0.05, 0.5)
        
        # Plot
        ax.scatter(x, y, z, c=colors, s=35, depthshade=False, 
                   marker='s', linewidth=0, edgecolors='none', antialiased=False)

    ax.view_init(elev=25, azim=45)
    
    img_stream = io.BytesIO()
    plt.savefig(img_stream, format='png', dpi=100, facecolor='black')
    plt.close(fig)
    img_stream.seek(0)
    return img_stream


def process_single_object(data_package):
    """
    Worker function: Calculates profile, normalizes, detects OUTLIERS, 
    and generates visualization.
    """
    (obj_id, cz, cy, cx, area, crop_r, crop_g, crop_mask) = data_package
    
    # 1. Calculate Raw Profile
    raw_r, raw_g = calculate_spherical_profile(crop_r, crop_g, crop_mask, int(obj_id))
    
    # 2. Normalize (Min = 1)
    def normalize_min_1(arr):
        min_val = np.min(arr)
        if min_val <= 0:
            if np.max(arr) == 0: return arr
            return arr 
        return arr / min_val

    norm_r = normalize_min_1(raw_r)
    norm_g = normalize_min_1(raw_g)
    
    # --- NEW: Outlier Detection ---
    # Criterion: If normalized max intensity > 30
    is_outlier = False
    if np.max(norm_r) > 30 or np.max(norm_g) > 30:
        is_outlier = True

    # 3. Classification
    size_class = "small" if area <= 90 else "large" 
    trend_class = trend_tools.analyze_profile_trend(norm_r, norm_g)
    
    # 4. Generate Images
    x_vals = np.arange(0, 11)
    # Add "(OUTLIER)" to plot title if detected
    plot_title = f"Object: {obj_id} (OUTLIER)" if is_outlier else f"Object: {obj_id}"
    
    # We need to update create_plot_image signature slightly or just pass ID
    # For now, let's keep the ID as the title carrier
    plot_bytes = create_plot_image(norm_r, norm_g, x_vals, plot_title).getvalue()
    
    thumb_bytes = create_dimetric_thumbnail(crop_r, crop_g).getvalue()
    
    return {
        "id": obj_id,
        "cz": cz, "cy": cy, "cx": cx,
        "area": area,
        "norm_r": norm_r,
        "norm_g": norm_g,
        "size_class": size_class,
        "trend_class": trend_class,
        "plot_bytes": plot_bytes,
        "thumb_bytes": thumb_bytes,
        "is_outlier": is_outlier  # Passing flag back to main
    }

# --- MAIN CONTROLLER ---

def main(red_path, green_path, mask_path, output_pptx):
    print(f"Loading 3D images from {os.path.dirname(red_path)}...")
    
    try:
        # --- 1. MEMMAP LOADING (Prevents RAM Saturation) ---
        print(f" Mapping Red Channel: {os.path.basename(red_path)}...")
        img_r = tifffile.memmap(red_path)
        
        print(f" Mapping Green Channel: {os.path.basename(green_path)}...")
        img_g = tifffile.memmap(green_path)
        
        print(f" Loading Mask: {os.path.basename(mask_path)}...")
        # Mask is usually smaller/compressible, so we load it fully for labeling
        mask = tifffile.imread(mask_path)
        
    except Exception as e:
        print(f"\nCRITICAL ERROR loading images: {e}")
        return

    # Handle extra dimensions
    if mask.ndim > 3: mask = mask.squeeze()

    # Labeling
    unique_vals = np.unique(mask)
    if len(unique_vals) <= 2 or mask.max() == 255:
        print("Labeling binary mask in 3D...")
        labeled_mask = measure.label(mask > 0)
    else:
        labeled_mask = mask

    props = measure.regionprops(labeled_mask)
    print(f"Found {len(props)} objects. Preparing tasks...")

    # Create Tasks
    tasks = []
    d_img, h_img, w_img = img_r.shape
    crop_rad = 10 
    
    for prop in props:
        cz, cy, cx = map(int, prop.centroid)
        
        # Calculate 3D bounds
        z1, z2 = max(0, cz - crop_rad), min(d_img, cz + crop_rad + 1)
        y1, y2 = max(0, cy - crop_rad), min(h_img, cy + crop_rad + 1)
        x1, x2 = max(0, cx - crop_rad), min(w_img, cx + crop_rad + 1)
        
        # Read Data (This reads from disk on-demand via memmap)
        crop_r = img_r[z1:z2, y1:y2, x1:x2]
        crop_g = img_g[z1:z2, y1:y2, x1:x2]
        crop_m = labeled_mask[z1:z2, y1:y2, x1:x2]
        
        tasks.append((
            prop.label, cz, cy, cx, prop.area,
            crop_r, crop_g, crop_m
        ))

    # --- PPT INITIALIZATION ---
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    
    groups = {} # For Summary Statistics

    print(f"Processing {len(tasks)} objects...")

    # --- PARALLEL PROCESSING & SLIDE GENERATION ---
    # FIX: The loop MUST be inside this 'with' block
    with concurrent.futures.ProcessPoolExecutor() as executor:
        
        # executor.map processes tasks in order and yields results as they finish
        results_generator = executor.map(process_single_object, tasks)
        
        for i, result in enumerate(results_generator):
            if result is None: continue
            
            # Progress Counter
            if (i + 1) % 10 == 0:
                print(f" Completed {i + 1}/{len(tasks)}...")

            # --- 1. HANDLE OUTLIERS ---
            if result['is_outlier']:
                print(f"  > Object {result['id']}: OUTLIER (Max > 30). Excluded from stats.")
            else:
                # Only add non-outliers to the Average Group
                key = (result['trend_class'], result['size_class'])
                if key not in groups: groups[key] = {'r': [], 'g': []}
                groups[key]['r'].append(result['norm_r'])
                groups[key]['g'].append(result['norm_g'])

            # --- 2. CREATE SLIDE (For ALL objects) ---
            slide = prs.slides.add_slide(blank_layout)
            
            # Add Graph
            plot_stream = io.BytesIO(result['plot_bytes'])
            slide.shapes.add_picture(plot_stream, left=Inches(0.5), top=Inches(1), width=Inches(7.5), height=Inches(5))
            
            # Add 3D Dimetric Thumbnail
            thumb_stream = io.BytesIO(result['thumb_bytes'])
            slide.shapes.add_picture(thumb_stream, left=Inches(8), top=Inches(0.5), width=Inches(3), height=Inches(3))
            
            # Add Text Box
            tx = slide.shapes.add_textbox(left=Inches(8), top=Inches(4), width=Inches(5), height=Inches(3))
            tf = tx.text_frame
            tf.word_wrap = True
            
            # Define Status Text
            status_line = "Analysis: OK"
            if result['is_outlier']:
                status_line = "STATUS: OUTLIER (Excluded)"

            lines = [
                f"Object ID: {result['id']}",
                f"Coords(z,y,x): ({result['cz']}, {result['cy']}, {result['cx']})",
                f"Volume(px): {result['area']}",
                f"Trend: {result['trend_class']}",
                f"Size: {result['size_class']}",
                status_line
            ]
            
            for line in lines:
                p = tf.add_paragraph()
                p.text = line
                p.font.size = Inches(0.25)
                
                # Bolding logic
                if "Trend" in line and result['trend_class'] != "No trend":
                    p.font.bold = True
                
                # Red text for Outliers
                if "OUTLIER" in line:
                    p.font.color.rgb = RGBColor(255, 0, 0)
                    p.font.bold = True

    # --- SUMMARY SLIDES ---
    print("Generating Summary Slides...")
    
    trend_types = ["Colocalized", "Around", "Anticolocalized", "No trend"]
    size_types = ["small", "large"] 
    
    for trend in trend_types:
        for size in size_types:
            key = (trend, size)
            
            if key in groups and len(groups[key]['r']) > 0:
                slide = prs.slides.add_slide(blank_layout)
                
                # Calculate Average of Normalized Profiles
                avg_r = np.mean(groups[key]['r'], axis=0)
                avg_g = np.mean(groups[key]['g'], axis=0)
                
                # Create Plot for Average
                plot_stream = create_plot_image(avg_r, avg_g, np.arange(11), "AVERAGE")
                slide.shapes.add_picture(plot_stream, Inches(0.5), Inches(1), height=Inches(5))
                
                # Info Text
                tx = slide.shapes.add_textbox(Inches(8), Inches(2.5), Inches(5), Inches(2))
                tf = tx.text_frame
                tf.text = f"Group Average\nTrend: {trend}\nSize: {size}\nCount (n): {len(groups[key]['r'])}"
                for p in tf.paragraphs:
                    p.font.size = Inches(0.25)
                    p.font.bold = True

    # --- SAVE OUTPUT ---
    while True:
        try:
            prs.save(output_pptx)
            print(f"Saved successfully: {output_pptx}")
            break
        except PermissionError:
            input(f"File {output_pptx} is OPEN. Close it and press ENTER...")


if __name__ == "__main__":
    # --- CONFIGURATION ---
    # Update these paths to point to your 3D TIFF files
    # Assuming input is now a 3D stack (Z, Y, X)
    
    base_path = r"/media/kyoung/Elements/20260101 PFKL-mCherry_EGFP-GAPDH - High Salt Conc - 24.1 degree - LLSM/Deconv/Cell2 - 180mM/3_NoGlucose60min"
    
    # Example Filenames (Ensure these exist as 3D Tiffs)
    # If your files are split (one file per z-slice), you need to stack them first or load as a sequence.
    # This code assumes Multipage TIFF (ImageJ default for 3D stacks)
    
    if os.path.exists(base_path):
        main(
            os.path.join(base_path, "Cell2_ex560em593_250mW_200ms_NoGlucose60min_d62_corrected.tif"), # Red 3D Stack
            os.path.join(base_path, "Cell2_ex488em514_230mW_200ms_NoGlucose60min_d16.tif"), # Green 3D Stack
            os.path.join(base_path, "Cell2_ex560em593_250mW_200ms_NoGlucose60min_d62_corrected-1.segmentation.tif"),      # Mask 3D Stack
            os.path.join(base_path, "Cell2_NoGlucose60min_3D_Report.pptx")
        )
    else:
        print("Please configure the 'base_path' in the script.")