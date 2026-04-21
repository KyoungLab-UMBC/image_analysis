import matplotlib
# Force Agg backend to prevent GUI errors
matplotlib.use('Agg') 
import matplotlib.pyplot as plt
# Plotting imports
from matplotlib.figure import Figure
from matplotlib.backends.backend_agg import FigureCanvasAgg

import numpy as np
import tifffile
from skimage import measure
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
import io
import warnings
import os
import re
import pandas as pd
from numba import jit
import concurrent.futures
import threading
render_lock = threading.Lock()

# --- CUSTOM MODULE IMPORTS ---
import util.trend_analysis as trend_tools
import util.render_3d as render_3d
import distance_analysis as dis_tool
import util.config as cfg
import util.region_separation as rs

# Set matplotlib backend to Agg for thread safety (no GUI)
plt.switch_backend('Agg')

# Suppress warnings
warnings.filterwarnings("ignore", category=RuntimeWarning)
warnings.filterwarnings("ignore", category=UserWarning)

# --- ACCELERATED WORKER FUNCTIONS (NUMBA) ---

@jit(nopython=True)
def calculate_spherical_profile(crop_r, crop_g, crop_mask, crop_roi_mask, obj_id):
    """
    Numba-accelerated function to calculate radial profile in 3D.
    Iterates over the 21x21x21 cube and bins pixels by spherical radius.
    """
    d, h, w = crop_r.shape
    cz, cy, cx = d // 2, h // 2, w // 2
    
    # Dynamically calculate the maximum radius based on the cube size
    max_r = 21
    sums_r = np.zeros(max_r)
    sums_g = np.zeros(max_r)
    counts = np.zeros(max_r)
    
    for z in range(d):
        for y in range(h):
            for x in range(w):
                # 1. ROI Mask Logic: Skip if the pixel is outside the ROI
                if crop_roi_mask[z, y, x] == 0:
                    continue
                
                # 2. Object Mask Logic: Include if background (0) or current object (obj_id)
                if crop_mask[z, y, x] != 0 and crop_mask[z, y, x] != obj_id:
                    continue
                
                # 3D Distance
                dist = np.sqrt((z - cz)**2 + (y - cy)**2 + (x - cx)**2)
                
                r_idx = int(dist)
                
                if r_idx < max_r:
                    sums_r[r_idx] += crop_r[z, y, x]
                    sums_g[r_idx] += crop_g[z, y, x]
                    counts[r_idx] += 1
                    
    means_r = np.zeros(max_r)
    means_g = np.zeros(max_r)
    
    for i in range(max_r):
        if counts[i] > 0:
            means_r[i] = sums_r[i] / counts[i]
            means_g[i] = sums_g[i] / counts[i]
            
    return means_r, means_g


def create_plot_image(red_data, green_data, obj_id, r_name, g_name):
    # Use the object-oriented API which is isolated and thread-safe
    fig = Figure(figsize=(6, 4))
    canvas = FigureCanvasAgg(fig)
    
    ax1 = fig.add_subplot(111)
    
    # Generate x-axis based on the length of the data
    data_length = len(red_data)
    x_axis = np.arange(data_length)
    
    ax1.set_xlabel('Radius (pixel)')
    ax1.set_xlim(0, data_length - 1)
    ax1.set_xticks(range(0, data_length, 1))
    
    # Left Y: Red
    ax1.set_ylabel(f'Norm. {r_name}', color='red')
    ax1.plot(x_axis, red_data, color='red', linewidth=1.5)
    ax1.tick_params(axis='y', labelcolor='red')

    # Right Y: Green
    ax2 = ax1.twinx() 
    ax2.set_ylabel(f'Norm. {g_name}', color='green')
    ax2.plot(x_axis, green_data, color='green', linewidth=1.5)
    ax2.tick_params(axis='y', labelcolor='green')

    if obj_id != "AVERAGE":
        ax1.set_title(f"Object: {obj_id}")
    else:
        ax1.set_title("Group Average (Normalized)")
    
    fig.tight_layout()
    
    img_stream = io.BytesIO()
    canvas.print_png(img_stream)  # Use canvas to print to stream directly
    img_stream.seek(0)
    
    return img_stream


def process_single_object(data_package):
    """
    Worker function: Calculates profile, normalizes, detects OUTLIERS, 
    and generates visualization matching 2D conventions.
    """
    (obj_id, cz, cy, cx, area, crop_r, crop_g, crop_m, crop_mask, crop_roi_mask, mito_mode, dist_val, r_name, g_name, generate_slides) = data_package
    
    # Object-specific mask
    obj_mask = (crop_mask == obj_id)
    
    # Raw intensities and Max inside
    r_in = np.max(crop_r[obj_mask]) if np.any(obj_mask) else 0
    g_in = np.max(crop_g[obj_mask]) if np.any(obj_mask) else 0
    raw_int_r = np.sum(crop_r[obj_mask])
    raw_int_g = np.sum(crop_g[obj_mask])

    # 1. Calculate Raw Profile
    raw_r, raw_g = calculate_spherical_profile(crop_r, crop_g, crop_mask, crop_roi_mask, int(obj_id))
    
    # 2. Normalize (Min = 1 approach for 3D to avoid heavy BG calculation)
    def normalize_min_1(arr):
        min_val = np.min(arr)
        if min_val <= 0:
            if np.max(arr) == 0: return arr
            return arr 
        return arr / min_val

    norm_r = normalize_min_1(raw_r)
    norm_g = normalize_min_1(raw_g)
    
    # 3. Classification
    size_class = "small glucosome" if area <= 90 else "large glucosome" 
    trend_class, r_value, sigma, pc_r, pc_g = trend_tools.analyze_profile_trend(norm_r, norm_g)
    
    # 4. Outlier Detection
    is_outlier = False
    if np.max(norm_r) > 30 or np.max(norm_g) > 30 or np.argmax(norm_r) not in [0, 1] or area < 16:
        is_outlier = True

    # 5. Distance Based Classification
    mito_status = "N/A"
    if mito_mode:
        if dist_val is not None and dist_val <= 275:
            mito_status = "Associated"
        else:
            mito_status = "Dissociated"

    # 6. Generate Images
    plot_bytes = b""
    thumb1_bytes = b""
    thumb2_bytes = b""
    
    if generate_slides:
        # Matplotlib is thread-safe here because we used the Object-Oriented API
        plot_bytes = create_plot_image(norm_r, norm_g, obj_id, r_name, g_name).getvalue()
        
        # VTK is NOT thread-safe. We must lock this specific step.
        with render_lock:
            thumb1_bytes, thumb2_bytes = render_3d.create_dimetric_thumbnails_VTK(crop_r, crop_g)
                
    return {
        "id": obj_id,
        "cz": cz, "cy": cy, "cx": cx,
        "area": area,
        "norm_r": norm_r,
        "norm_g": norm_g,
        "size_class": size_class,
        "trend_class": trend_class,
        "r_value": r_value,
        "sigma": sigma,
        "pc_r": pc_r,   
        "pc_g": pc_g,   
        "r_in": r_in,           
        "g_in": g_in,           
        "raw_int_r": raw_int_r, 
        "raw_int_g": raw_int_g,
        "dist_val": dist_val, 
        "mito_status": mito_status, 
        "is_outlier": is_outlier,
        "plot_bytes": plot_bytes,
        "thumb1_bytes": thumb1_bytes,
        "thumb2_bytes": thumb2_bytes
    }

# --- MAIN CONTROLLER ---

def processing_cell_3d(user_path, red_path, green_path, cell_mask_path, mito_mode, r_name, g_name, generate_slides):
    img_m = None
    dist_map = {}
    mito_path = None
    roi_mask = tifffile.imread(cell_mask_path) > 0
    roi_name = f"{cell_mask_path.stem}"
    if mito_mode:
        mito_path = user_path / f"{cfg.MITO_RAW.stem}_{roi_name}_mask.tif"
    # Construct paths using stem and roi_name (cell_num)
    r_mask_path = user_path / f"{cfg.R_RAW.stem}_{roi_name}_mask.tif"
    r_mask = tifffile.imread(r_mask_path)
    output_pptx = user_path / f"{cfg.R_RAW.stem}_{roi_name}.pptx"

    try:
        # --- 1. MEMMAP LOADING ---
        print(f" Mapping Red Channel...")
        img_r = tifffile.imread(red_path)
        
        print(f" Mapping Green Channel...")
        img_g = tifffile.imread(green_path)
        
        print(f" Loading Source Mask...")
        r_mask = tifffile.imread(r_mask_path)
        
    except Exception as e:
        print(f"\nCRITICAL ERROR loading images: {e}")
        return

    if r_mask.ndim > 3: r_mask = r_mask.squeeze()

    # Labeling
    unique_vals = np.unique(r_mask)
    if len(unique_vals) <= 2 or r_mask.max() == 255:
        print("Labeling binary mask in 3D...")
        labeled_mask = measure.label(r_mask > 0)
    else:
        labeled_mask = r_mask

    # Distance Analysis
    if mito_mode and mito_path and mito_path.exists():
        print("Loading Mito Mask and calculating 3D distances...")
        # 1. Get objects from Mask (Source)
        pfkl_objs, _ = dis_tool.get_object_data(r_mask_path, is_2d_mode=False)
        # 2. Get objects from Mito Mask (Target Cloud)
        _, mito_cloud = dis_tool.get_object_data(mito_path, is_2d_mode=False)
        # 3. Calculate Distances (Assuming standard LLSM Z-step ~250nm or keep 107 if isotropic)
        dists = dis_tool.analyze_distances(pfkl_objs, mito_cloud, is_2d_mode=False, s_xy=107, s_z=250) 
        
        for obj, d in zip(pfkl_objs, dists):
            dist_map[obj['label']] = d

    props = measure.regionprops(labeled_mask)
    print(f"Found {len(props)} objects. Preparing tasks...")

    # Calculate Total Intensity for Excel
    # For 3D, calculating total sum over entire memmap might be slow, doing it over mask
    total_roi_int_r = np.sum(img_r[roi_mask])
    total_roi_int_g = np.sum(img_g[roi_mask])
    if total_roi_int_r == 0: total_roi_int_r = 1
    if total_roi_int_g == 0: total_roi_int_g = 1

    # Delegate logic to region_separation.py
    r_bg, g_bg, mask_hh, mask_lh, mask_hl, mask_ll = rs.separate_regions_and_bg(img_r, img_g, roi_mask)
    
    # Save the 4 quadrant masks to a zip file
    regions_zip_path = user_path / f"{cfg.R_RAW.stem}_{roi_name}_regions.zip"
    rs.save_masks_to_tifs(regions_zip_path, mask_hh, mask_lh, mask_hl, mask_ll)

    def get_out_val(bg_img, mask):
        vals = bg_img[mask].flatten()
        # Sample rate 1/1000
        sampled = vals[::1000]
        if len(sampled) == 0: return 1.0       
        # Sort values in ascending order (lowest to highest)
        sampled = np.sort(sampled)        
        # Mean of first 10% (lowest values in the sampled group)
        idx = max(1, int(len(sampled) * 0.1))
        return np.mean(sampled[:idx])
    
    hh_r_out = get_out_val(r_bg, mask_hh)
    hh_g_out = get_out_val(g_bg, mask_hh)
    
    lh_r_out = get_out_val(r_bg, mask_lh)
    lh_g_out = get_out_val(g_bg, mask_lh)
    
    hl_r_out = get_out_val(r_bg, mask_hl)
    hl_g_out = get_out_val(g_bg, mask_hl)
    
    ll_r_out = get_out_val(r_bg, mask_ll)
    ll_g_out = get_out_val(g_bg, mask_ll)

    # Create Tasks
    tasks = []
    d_img, h_img, w_img = img_r.shape
    crop_rad = 20 
    
    for prop in props:
        cz, cy, cx = map(int, prop.centroid)
        
        z1, z2 = max(0, cz - crop_rad), min(d_img, cz + crop_rad + 1)
        y1, y2 = max(0, cy - crop_rad), min(h_img, cy + crop_rad + 1)
        x1, x2 = max(0, cx - crop_rad), min(w_img, cx + crop_rad + 1)
        
        crop_r = img_r[z1:z2, y1:y2, x1:x2]
        crop_g = img_g[z1:z2, y1:y2, x1:x2]
        crop_m = None
        if mito_mode and img_m is not None:
            crop_m = img_m[z1:z2, y1:y2, x1:x2]

        crop_mask = labeled_mask[z1:z2, y1:y2, x1:x2]
        crop_roi_mask = roi_mask[z1:z2, y1:y2, x1:x2]

        dist_val = dist_map.get(prop.label, 99999) if mito_mode else None
        
        tasks.append((
            prop.label, cz, cy, cx, prop.area,
            crop_r, crop_g, crop_m, crop_mask, crop_roi_mask,
            mito_mode, dist_val, r_name, g_name, generate_slides
        ))

    results = []
    max_workers = min(32, os.cpu_count())
    print(f"Processing {len(tasks)} objects with ThreadPool ({max_workers} workers)...")

    with concurrent.futures.ThreadPoolExecutor(max_workers=max_workers) as executor:
        for i, res in enumerate(executor.map(process_single_object, tasks)):
            results.append(res)
            if (i + 1) % 50 == 0:
                print(f" Completed {i + 1}/{len(tasks)}...")

    # --- PPT GENERATION ---
    groups = {}
    ignore_trend = (cfg.G_PATH == cfg.MITO_RAW)

    if generate_slides:
        print("Generating PowerPoint Report...")
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank_layout = prs.slide_layouts[6]

    for res in results:
        if mito_mode:
            if ignore_trend:
                k = (res['size_class'], res['mito_status'])
            else:
                k = (res['trend_class'], res['size_class'], res['mito_status'])
        else:
            k = (res['trend_class'], res['size_class'])

        cz, cy, cx = int(res['cz']), int(res['cy']), int(res['cx'])
        if mask_hh[cz, cy, cx]:
            r_out, g_out = hh_r_out, hh_g_out
        elif mask_lh[cz, cy, cx]:
            r_out, g_out = lh_r_out, lh_g_out
        elif mask_hl[cz, cy, cx]:
            r_out, g_out = hl_r_out, hl_g_out
        else: # mask_ll
            r_out, g_out = ll_r_out, ll_g_out
            
        res['r_out'] = r_out
        res['g_out'] = g_out

        res['pc_r_io'] = res['r_in'] / r_out if r_out != 0 else 0
        res['pc_g_io'] = res['g_in'] / g_out if g_out != 0 else 0
        res['pc_ratio_io'] = res['pc_g_io'] / res['pc_r_io'] if res['pc_r_io'] != 0 else 0
        
        if not res['is_outlier']:
            if k not in groups: groups[k] = {'r': [], 'g': []}
            groups[k]['r'].append(res['norm_r'])
            groups[k]['g'].append(res['norm_g'])
        
        if generate_slides:
            slide = prs.slides.add_slide(blank_layout)
            
            # Graphs and Thumbs
            slide.shapes.add_picture(io.BytesIO(res['plot_bytes']), left=Inches(0.5), top=Inches(1), width=Inches(7.5), height=Inches(5))
            slide.shapes.add_picture(io.BytesIO(res['thumb1_bytes']), left=Inches(8), top=Inches(0.5), width=Inches(2.5), height=Inches(2.5))
            slide.shapes.add_picture(io.BytesIO(res['thumb2_bytes']), left=Inches(10.5), top=Inches(0.5), width=Inches(2.5), height=Inches(2.5))
            
            tx = slide.shapes.add_textbox(left=Inches(8), top=Inches(4), width=Inches(5), height=Inches(3))
            tf = tx.text_frame
            tf.word_wrap = True

            # Legends
            if mito_mode:
                p = tf.add_paragraph(); p.text = "Mitochondria"; p.font.color.rgb = RGBColor(255, 0, 0); p.font.size = Inches(0.25); p.font.bold = True
                p = tf.add_paragraph(); p.text = r_name; p.font.color.rgb = RGBColor(0, 255, 0); p.font.size = Inches(0.25); p.font.bold = True
                p = tf.add_paragraph(); p.text = g_name; p.font.color.rgb = RGBColor(0, 0, 255); p.font.size = Inches(0.25); p.font.bold = True
            else:
                p = tf.add_paragraph(); p.text = r_name; p.font.color.rgb = RGBColor(255, 0, 0); p.font.size = Inches(0.25); p.font.bold = True
                p = tf.add_paragraph(); p.text = g_name; p.font.color.rgb = RGBColor(0, 255, 0); p.font.size = Inches(0.25); p.font.bold = True

            lines = [
                f"Object ID: {res['id']}",
                f"Coords(x,y,z): ({res['cx']}, {res['cy']}, {res['cz']})",
                f"Volume(voxels): {res['area']}",
                f"Sigma: {res['sigma']:.2f}",
                f"Trend: {res['trend_class']}",
                f"Pearson coefficient: {res['r_value']:.2f}"
            ]
            if mito_mode: lines.append(f"Mito: {res['mito_status']}")
            if res['is_outlier']: lines.append("!!! OUTLIER - EXCLUDED !!!")

            for line in lines:
                p = tf.add_paragraph()
                p.text = line
                p.font.size = Inches(0.25)
                if "Trend" in line and res['trend_class'] != "No trend": p.font.bold = True
                if "Mito" in line and "Associated" in line and "Not" not in line: p.font.bold = True

    # Categories Setup
    trend_types = ["Colocalized", "Around", "Anticolocalized", "No trend"] if not ignore_trend else [None]
    size_types = ["small glucosome", "large glucosome"]
    mito_types = ["Associated", "Dissociated"] if mito_mode else [None]

    categories = []
    if mito_mode:
        if ignore_trend:
            for size in size_types:
                    for mito in mito_types: categories.append((size, mito))
        else:
            for trend in trend_types:
                for size in size_types:
                    for mito in mito_types: categories.append((trend, size, mito))
    else:
        for trend in trend_types:
            for size in size_types: categories.append((trend, size))

    # --- 1. SUMMARY SLIDES ---
    if generate_slides:
        print("Generating Summary Slides...")
        for cat_tuple in categories:
            if mito_mode:
                if ignore_trend:
                    size, mito = cat_tuple
                    key = (size, mito)
                    title_str = f"Group Average\n{size}\n{mito}"
                else:
                    trend, size, mito = cat_tuple
                    key = (trend, size, mito)
                    title_str = f"Group Average\n{trend}\n{size}\n{mito}"
            else:
                trend, size = cat_tuple
                key = (trend, size)
                title_str = f"Group Average\n{trend}\n{size}"

            if key in groups and len(groups[key]['r']) > 0:
                count = len(groups[key]['r'])
                title_str += f"\n(n={count})"
                slide = prs.slides.add_slide(blank_layout)
                
                avg_r = np.mean(groups[key]['r'], axis=0)
                avg_g = np.mean(groups[key]['g'], axis=0)
                
                plot_stream = create_plot_image(avg_r, avg_g, "AVERAGE", r_name, g_name)
                slide.shapes.add_picture(plot_stream, Inches(0.5), Inches(1), height=Inches(5))
                
                tx = slide.shapes.add_textbox(Inches(8), Inches(2.5), Inches(5), Inches(2))
                tx.text_frame.text = title_str
                for p in tx.text_frame.paragraphs:
                    p.font.bold = True
                    p.font.size = Inches(0.25)

        while True:
            try:
                prs.save(output_pptx)
                print(f"Saved PowerPoint: {output_pptx}")
                break
            except PermissionError:
                input(f"File {output_pptx} is OPEN. Close it and press ENTER...")

    # --- 2. EXCEL REPORT ---
    print("Generating Excel Report...")
    output_xlsx = output_pptx.with_suffix(".xlsx")
    
    roi_area = np.sum(roi_mask)
    if roi_area == 0: roi_area = 1
    mean_roi_int_r = total_roi_int_r / roi_area
    mean_roi_int_g = total_roi_int_g / roi_area
    glucosome_num = len(results)


    folder_name = user_path.name
    prefix_num = 0
    match_prefix = re.match(r'^(\d+)_', folder_name)
    if match_prefix: prefix_num = int(match_prefix.group(1))
    time_point = folder_name if prefix_num != 0 else prefix_num

    try:
        with pd.ExcelWriter(output_xlsx, engine='openpyxl') as writer:
            
            path_str = str(user_path) 
            date_match = re.search(r'(\d{8})', path_str)
            date_part = date_match.group(1) if date_match else "00000000"
            cell_match = re.search(r'[Cc]ell(\d+)', path_str)
            cell_part = cell_match.group(1) if cell_match else "0"
            image_id = f"{date_part}-{cell_part}-{time_point}"

            summary_main = {
                "Image Name": image_id, "Area/px": roi_area, "Glucosome num": glucosome_num,
                f"Mean {r_name} int": mean_roi_int_r, f"Total {r_name} int": total_roi_int_r,
                f"Mean {g_name} int": mean_roi_int_g, f"Total {g_name} int": total_roi_int_g
            }

            summary_num = {"Image Name": image_id}
            summary_ratio_count = {"Image Name": image_id}
            summary_ratio_area = {"Image Name": image_id}
            summary_ratio_r = {"Image Name": image_id}
            summary_ratio_g = {"Image Name": image_id}
            
            valid_objects = [r for r in results if not r['is_outlier']]
            total_valid_count = len(valid_objects) if len(valid_objects) > 0 else 1

            for cat_tuple in categories:
                if mito_mode:
                    if ignore_trend:
                        size, mito = cat_tuple
                        group_items = [r for r in valid_objects if r['size_class'] == size and r['mito_status'] == mito]
                        group_name = f"{size} {mito}"
                    else:
                        trend, size, mito = cat_tuple
                        group_items = [r for r in valid_objects if r['trend_class'] == trend and r['size_class'] == size and r['mito_status'] == mito]
                        group_name = f"{trend} {size} {mito}"
                else:
                    trend, size = cat_tuple
                    group_items = [r for r in valid_objects if r['trend_class'] == trend and r['size_class'] == size]
                    group_name = f"{trend} {size}"
                
                count = len(group_items)
                summary_num[group_name] = count
                summary_ratio_count[group_name] = count / total_valid_count
                summary_ratio_area[group_name] = sum(item['area'] for item in group_items) / roi_area
                summary_ratio_r[group_name] = sum(item['raw_int_r'] for item in group_items) / total_roi_int_r
                summary_ratio_g[group_name] = sum(item['raw_int_g'] for item in group_items) / total_roi_int_g

            pd.DataFrame([summary_main]).to_excel(writer, sheet_name="Summary", index=False, startcol=prefix_num * 8)
            
            if mito_mode:
                if ignore_trend:
                    start_off = prefix_num * 6
                else:
                    start_off = prefix_num * 18
            else:
                start_off = prefix_num * 10
            
            pd.DataFrame([summary_num]).to_excel(writer, sheet_name="Summary_num", index=False, startcol=start_off)
            pd.DataFrame([summary_ratio_count]).to_excel(writer, sheet_name="Summary_ratio", index=False, startcol=start_off)
            pd.DataFrame([summary_ratio_area]).to_excel(writer, sheet_name="Summary_area_ratio", index=False, startcol=start_off)
            pd.DataFrame([summary_ratio_r]).to_excel(writer, sheet_name=f"Summary_{r_name}_ratio", index=False, startcol=start_off)
            pd.DataFrame([summary_ratio_g]).to_excel(writer, sheet_name=f"Summary_{g_name}_ratio", index=False, startcol=start_off)

            for cat_tuple in categories:
                if mito_mode:
                    if ignore_trend:
                        size, mito = cat_tuple
                        subset = [r for r in results if r['size_class'] == size and r['mito_status'] == mito and not r['is_outlier']]
                        sheet_name = f"{'Sm' if 'small' in size else 'Lg'} {'Att' if mito == 'Associated' else 'NotAtt'}"
                    else:
                        trend, size, mito = cat_tuple
                        subset = [r for r in results if r['trend_class'] == trend and r['size_class'] == size and r['mito_status'] == mito and not r['is_outlier']]
                        sheet_name = f"{trend[:5]} {'Sm' if 'small' in size else 'Lg'} {'Att' if mito == 'Associated' else 'NotAtt'}"
                else:
                    trend, size = cat_tuple
                    subset = [r for r in results if r['trend_class'] == trend and r['size_class'] == size and not r['is_outlier']]
                    sheet_name = f"{trend} {size}"[:31]
                
                if not subset: continue

                # Get the dynamic length from the first item in the subset
                data_len = len(subset[0]['norm_g'])
                
                header_row = [f"Normalized {g_name} {time_point}"] + list(range(data_len)) + ["", f"Normalized {r_name} {time_point}"] + list(range(data_len))
                data_rows = [header_row]
                for item in subset:
                    data_rows.append([item['id']] + list(item['norm_g']) + ["", item['id']] + list(item['norm_r']))
                
                pd.DataFrame(data_rows).to_excel(writer, sheet_name=sheet_name, index=False, header=False, startcol=prefix_num * 46)

                pc_sheet_name = f"{sheet_name}_pc"[:31]
                pc_headers = [
                    f"Obj ID {time_point}", "Sigma", "Volume (vx)", 
                    f"PC {r_name} Radius", f"PC {g_name} Radius", "PC Radius Ratio", 
                    f"{r_name} in", f"{r_name} out", f"{r_name} PC IO",
                    f"{g_name} in", f"{g_name} out", f"{g_name} PC IO",
                    "PC IO ratio (G/R)", "Mito Dist"
                ]
                pc_rows = []
                for item in subset:
                    ratio = (item['pc_g'] / item['pc_r']) if item['pc_r'] else 0
                    pc_rows.append([
                        item['id'], 
                        item['sigma'], 
                        item['area'], 
                        item['pc_r'], 
                        item['pc_g'], 
                        ratio,
                        item.get('r_in', ''), 
                        item.get('r_out', ''),
                        item.get('pc_r_io', ''),
                        item.get('g_in', ''),
                        item.get('g_out', ''),
                        item.get('pc_g_io', ''),
                        item.get('pc_ratio_io', ''),
                        item['dist_val'] if mito_mode else ""
                    ])
                pd.DataFrame(pc_rows, columns=pc_headers).to_excel(writer, sheet_name=pc_sheet_name, index=False, startcol=prefix_num * 15)

        print(f"Saved Excel: {output_xlsx}")
        
    except Exception as e:
        print(f"Failed to save Excel: {e}")

def main():    
    if not cfg.R_PATH.exists():
        print(f"Path not found: {cfg.R_PATH}")
    else:
        print(f"Running 3D analysis for: {cfg.R_NAME} vs {cfg.G_NAME}")
        # Execute the 3D process utilizing the config inputs
        # Assuming the config points to the 3D memmaps instead of 2D arrays 
        # (You can swap cfg.R_PATH to direct strings if your config is split)
        processing_cell_3d(
            user_path=cfg.USER_PATH,
            red_path=cfg.R_PATH,
            green_path=cfg.G_PATH,
            cell_mask_path=cfg.CELL_MASK,
            mito_mode=cfg.MITO_ANALYSIS,
            r_name=cfg.R_NAME,
            g_name=cfg.G_NAME,
            generate_slides=cfg.GENERATE_SLIDES
        )

if __name__ == "__main__":
    main()