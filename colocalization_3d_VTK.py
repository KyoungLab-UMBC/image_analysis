import numpy as np
import matplotlib.pyplot as plt
import tifffile
from skimage import measure
from pptx import Presentation
from pptx.util import Inches
import io
import warnings
import vtk
from vtk.util import numpy_support
import os
from scipy import signal
from numba import jit
import concurrent.futures
from pptx.dml.color import RGBColor
import trend_analysis as trend_tools

# Set matplotlib backend to Agg for thread safety (no GUI)
plt.switch_backend('Agg')

# Suppress warnings
warnings.filterwarnings("ignore", category=RuntimeWarning)
warnings.filterwarnings("ignore", category=UserWarning)

# --- ACCELERATED WORKER FUNCTIONS (NUMBA) ---

@jit(nopython=True)
def calculate_spherical_profile(crop_r, crop_g, crop_mask, obj_id):
    """
    Numba-accelerated function to calculate radial profile in 3D.
    Iterates over the 21x21x21 cube and bins pixels by spherical radius.
    """
    d, h, w = crop_r.shape
    cz, cy, cx = d // 2, h // 2, w // 2
    
    # Radius 0 to 10 (11 bins)
    max_r = 11
    sums_r = np.zeros(max_r)
    sums_g = np.zeros(max_r)
    counts = np.zeros(max_r)
    
    for z in range(d):
        for y in range(h):
            for x in range(w):
                # Mask Logic: Include if background (0) or current object (obj_id)
                # Exclude if it belongs to a neighbor object
                if crop_mask[z, y, x] != 0 and crop_mask[z, y, x] != obj_id:
                    continue
                
                # 3D Distance
                dist = np.sqrt((z - cz)**2 + (y - cy)**2 + (x - cx)**2)
                
                # Floor to get bin index (0.5 -> 0, 1.9 -> 1)
                r_idx = int(dist)
                
                if r_idx < max_r:
                    sums_r[r_idx] += crop_r[z, y, x]
                    sums_g[r_idx] += crop_g[z, y, x]
                    counts[r_idx] += 1
                    
    means_r = np.zeros(max_r)
    means_g = np.zeros(max_r)
    
    for i in range(max_r):
        if counts[i] > 0:
            means_r[i] = sums_r[i] / counts[i]
            means_g[i] = sums_g[i] / counts[i]
            
    return means_r, means_g


def create_plot_image(red_data, green_data, x_axis, obj_id):
    fig, ax1 = plt.subplots(figsize=(6, 4))
    
    ax1.set_xlabel('Radius (pixel)')
    ax1.set_xlim(0, 10)
    ax1.set_xticks(range(0, 11, 1))
    
    # Left Y: Red
    ax1.set_ylabel('Norm. Intensity (Red)', color='red')
    # Use standard scaling or dynamic based on data
    ax1.plot(x_axis, red_data, color='red', linewidth=1.5)
    ax1.tick_params(axis='y', labelcolor='red')

    # Right Y: Green
    ax2 = ax1.twinx() 
    ax2.set_ylabel('Norm. Intensity (Green)', color='green')
    ax2.plot(x_axis, green_data, color='green', linewidth=1.5)
    ax2.tick_params(axis='y', labelcolor='green')

    if obj_id != "AVERAGE":
        plt.title(f"Object: {obj_id}")
    else:
        plt.title("Group Average (Normalized)")
    
    fig.tight_layout()
    img_stream = io.BytesIO()
    plt.savefig(img_stream, format='png', dpi=100)
    plt.close(fig)
    img_stream.seek(0)
    return img_stream

def create_dimetric_thumbnails_rgba(crop_r, crop_g):
    """
    Creates TWO 3D Volume Renderings (View 1 and View 2 rotated 90 deg).
    Mixes Red and Green channels into a single RGBA texture.
    
    UPDATES:
    1. Box Gradient: Edges fade to 0.2 (Chebyshev distance).
    2. Green Transparency: Green contributes only 30% to opacity to prevent blocking.
    """
    
    # --- 1. PRE-PROCESSING & DYNAMIC RANGE ---
    s_start, s_end = 4, 17
    sub_r = crop_r[s_start:s_end, s_start:s_end, s_start:s_end].astype(np.float32)
    sub_g = crop_g[s_start:s_end, s_start:s_end, s_start:s_end].astype(np.float32)
    
    d, h, w = sub_r.shape
    cz, cy, cx = d // 2, h // 2, w // 2

    # --- Red Dynamic Range: Min to Center Pixel ---
    min_r = np.min(sub_r)
    max_r = sub_r[cz, cy, cx] # Center pixel sets the Max
    
    if max_r <= min_r: 
        max_r = np.max(sub_r) 
        if max_r <= min_r: max_r = min_r + 1e-5

    norm_r = (sub_r - min_r) / (max_r - min_r)
    norm_r = np.clip(norm_r, 0.0, 1.0) 

    # --- Green Dynamic Range: Min to Max ---
    min_g = np.min(sub_g)
    max_g = np.max(sub_g)
    if max_g <= min_g: max_g = min_g + 1e-5
    
    norm_g = (sub_g - min_g) / (max_g - min_g)
    norm_g = np.clip(norm_g, 0.0, 1.0)

    # --- 2. MIX TO RGBA ---
    rgba_vol = np.zeros((d, h, w, 4), dtype=np.uint8)
    
    # Colors remain full strength (0-255)
    rgba_vol[..., 0] = (norm_r * 255).astype(np.uint8) # R
    rgba_vol[..., 1] = (norm_g * 255).astype(np.uint8) # G
    rgba_vol[..., 2] = 0    # B
    
    # --- 3. CALCULATE SMART ALPHA ---
    # A. Box Distance (Chebyshev) - Handles Cube Shape
    zz, yy, xx = np.indices((d, h, w))
    box_dist = np.maximum(np.abs(zz - cz), np.maximum(np.abs(yy - cy), np.abs(xx - cx)))
    
    max_dist = d // 2
    if max_dist == 0: max_dist = 1
    
    # Gradient: 1.0 at center -> 0.2 at box edges
    spatial_scale = 1.0 - (box_dist / max_dist) * 0.8
    spatial_scale = np.clip(spatial_scale, 0.2, 1.0)
    
    # B. Opacity Mixing (The Fix for Blocking)
    # Green contributes significantly LESS to opacity (0.3) than Red (1.0).
    # This makes Green "ghostly" and Red "solid".
    base_opacity = np.maximum(norm_r, norm_g * 0.3)
    
    # Final Alpha = Weighted Intensity * Spatial Gradient
    alpha_float = base_opacity * spatial_scale
    alpha_u8 = (alpha_float * 255).astype(np.uint8)
    
    rgba_vol[..., 3] = alpha_u8

    # --- 4. VTK RENDERING ---
    data_flat = rgba_vol.reshape(-1, 4)
    vtk_data = numpy_support.numpy_to_vtk(num_array=data_flat, deep=True, array_type=vtk.VTK_UNSIGNED_CHAR)
    
    img = vtk.vtkImageData()
    img.SetDimensions(w, h, d)
    img.GetPointData().SetScalars(vtk_data)
    
    mapper = vtk.vtkSmartVolumeMapper()
    mapper.SetInputData(img)
    
    vol_prop = vtk.vtkVolumeProperty()
    vol_prop.ShadeOn()
    vol_prop.SetInterpolationTypeToLinear()
    vol_prop.IndependentComponentsOff() 
    
    # Alpha 0-255 maps to Opacity 0.0-1.0
    opacity_func = vtk.vtkPiecewiseFunction()
    opacity_func.AddPoint(0, 0.0)
    opacity_func.AddPoint(255, 1.0)
    vol_prop.SetScalarOpacity(opacity_func)
    
    vol_prop.SetAmbient(0.4)
    vol_prop.SetDiffuse(0.6)
    vol_prop.SetSpecular(0.2)
    
    volume = vtk.vtkVolume()
    volume.SetMapper(mapper)
    volume.SetProperty(vol_prop)
    
    renderer = vtk.vtkRenderer()
    renderer.SetBackground(0, 0, 0)
    renderer.AddVolume(volume)
    
    # Outline
    outline = vtk.vtkOutlineFilter()
    outline.SetInputData(img)
    mapper_outline = vtk.vtkPolyDataMapper()
    mapper_outline.SetInputConnection(outline.GetOutputPort())
    actor_outline = vtk.vtkActor()
    actor_outline.SetMapper(mapper_outline)
    actor_outline.GetProperty().SetColor(1, 1, 1)
    actor_outline.GetProperty().SetLineWidth(1.0)
    renderer.AddActor(actor_outline)

    # View 1
    render_window = vtk.vtkRenderWindow()
    render_window.SetOffScreenRendering(1)
    render_window.AddRenderer(renderer)
    render_window.SetSize(400, 400)
    
    camera = renderer.GetActiveCamera()
    camera.Elevation(20)
    camera.Azimuth(30)
    renderer.ResetCamera()
    camera.Zoom(1.3)
    
    render_window.Render()
    
    w2if = vtk.vtkWindowToImageFilter()
    w2if.SetInput(render_window)
    w2if.SetInputBufferTypeToRGB()
    w2if.ReadFrontBufferOff()
    w2if.Update()
    
    writer = vtk.vtkPNGWriter()
    writer.SetInputConnection(w2if.GetOutputPort())
    writer.WriteToMemoryOn()
    writer.Write()
    view1_bytes = bytes(writer.GetResult())

    # View 2 (Rotate 90)
    camera.Azimuth(90)
    render_window.Render()
    w2if.Modified()
    w2if.Update()
    writer.Write()
    view2_bytes = bytes(writer.GetResult())
    
    del volume, mapper, img, renderer, render_window, writer, w2if
    
    return view1_bytes, view2_bytes


def process_single_object(data_package):
    """
    Worker function: Calculates profile, normalizes, detects OUTLIERS, 
    and generates visualization.
    """
    (obj_id, cz, cy, cx, area, crop_r, crop_g, crop_mask) = data_package
    
    # 1. Calculate Raw Profile
    raw_r, raw_g = calculate_spherical_profile(crop_r, crop_g, crop_mask, int(obj_id))
    
    # 2. Normalize (Min = 1)
    def normalize_min_1(arr):
        min_val = np.min(arr)
        if min_val <= 0:
            if np.max(arr) == 0: return arr
            return arr 
        return arr / min_val

    norm_r = normalize_min_1(raw_r)
    norm_g = normalize_min_1(raw_g)
    
    # --- Outlier Detection ---
    is_outlier = False
    if np.max(norm_r) > 30 or np.max(norm_g) > 30 or not (np.argmax(norm_r) in [0, 1]):
        is_outlier = True

    # 3. Classification
    size_class = "small" if area <= 90 else "large" 
    trend_class = trend_tools.analyze_profile_trend(norm_r, norm_g)
    
    # 4. Generate Images
    x_vals = np.arange(0, 11)
    
    plot_bytes = create_plot_image(norm_r, norm_g, x_vals, obj_id).getvalue()
    
    # NEW: Generate Two Views with RGBA mixing and BOX Alpha Gradient
    thumb1_bytes, thumb2_bytes = create_dimetric_thumbnails_rgba(crop_r, crop_g)
    
    return {
        "id": obj_id,
        "cz": cz, "cy": cy, "cx": cx,
        "area": area,
        "norm_r": norm_r,
        "norm_g": norm_g,
        "size_class": size_class,
        "trend_class": trend_class,
        "plot_bytes": plot_bytes,
        "thumb1_bytes": thumb1_bytes,
        "thumb2_bytes": thumb2_bytes,
        "is_outlier": is_outlier
    }

# --- MAIN CONTROLLER ---

def main(red_path, green_path, mask_path, output_pptx):
    print(f"Loading 3D images from {os.path.dirname(red_path)}...")
    
    try:
        # --- 1. MEMMAP LOADING ---
        print(f" Mapping Red Channel: {os.path.basename(red_path)}...")
        img_r = tifffile.memmap(red_path)
        
        print(f" Mapping Green Channel: {os.path.basename(green_path)}...")
        img_g = tifffile.memmap(green_path)
        
        print(f" Loading Mask: {os.path.basename(mask_path)}...")
        mask = tifffile.imread(mask_path)
        
    except Exception as e:
        print(f"\nCRITICAL ERROR loading images: {e}")
        return

    if mask.ndim > 3: mask = mask.squeeze()

    # Labeling
    unique_vals = np.unique(mask)
    if len(unique_vals) <= 2 or mask.max() == 255:
        print("Labeling binary mask in 3D...")
        labeled_mask = measure.label(mask > 0)
    else:
        labeled_mask = mask

    props = measure.regionprops(labeled_mask)
    print(f"Found {len(props)} objects. Preparing tasks...")

    # Create Tasks
    tasks = []
    d_img, h_img, w_img = img_r.shape
    crop_rad = 10 
    
    for prop in props:
        cz, cy, cx = map(int, prop.centroid)
        
        z1, z2 = max(0, cz - crop_rad), min(d_img, cz + crop_rad + 1)
        y1, y2 = max(0, cy - crop_rad), min(h_img, cy + crop_rad + 1)
        x1, x2 = max(0, cx - crop_rad), min(w_img, cx + crop_rad + 1)
        
        crop_r = img_r[z1:z2, y1:y2, x1:x2]
        crop_g = img_g[z1:z2, y1:y2, x1:x2]
        crop_m = labeled_mask[z1:z2, y1:y2, x1:x2]
        
        tasks.append((
            prop.label, cz, cy, cx, prop.area,
            crop_r, crop_g, crop_m
        ))

    # --- PPT INITIALIZATION ---
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    
    groups = {} 

    print(f"Processing {len(tasks)} objects...")

    # --- PARALLEL PROCESSING & SLIDE GENERATION ---
    with concurrent.futures.ProcessPoolExecutor() as executor:
        
        results_generator = executor.map(process_single_object, tasks)
        
        for i, result in enumerate(results_generator):
            if result is None: continue
            
            if (i + 1) % 50 == 0:
                print(f" Completed {i + 1}/{len(tasks)}...")

            else:
                key = (result['trend_class'], result['size_class'])
                if key not in groups: groups[key] = {'r': [], 'g': []}
                groups[key]['r'].append(result['norm_r'])
                groups[key]['g'].append(result['norm_g'])

            # --- CREATE SLIDE ---
            slide = prs.slides.add_slide(blank_layout)
            
            # Add Graph
            plot_stream = io.BytesIO(result['plot_bytes'])
            slide.shapes.add_picture(plot_stream, left=Inches(0.5), top=Inches(1), width=Inches(7.5), height=Inches(5))
            
            # Add 3D View 1 (Standard)
            thumb1_stream = io.BytesIO(result['thumb1_bytes'])
            slide.shapes.add_picture(thumb1_stream, left=Inches(8), top=Inches(0.5), width=Inches(2.5), height=Inches(2.5))

            # Add 3D View 2 (Rotated 90 deg right)
            thumb2_stream = io.BytesIO(result['thumb2_bytes'])
            slide.shapes.add_picture(thumb2_stream, left=Inches(10.5), top=Inches(0.5), width=Inches(2.5), height=Inches(2.5))
            
            # Add Text Box
            tx = slide.shapes.add_textbox(left=Inches(8), top=Inches(4), width=Inches(5), height=Inches(3))
            tf = tx.text_frame
            tf.word_wrap = True
            
            status_line = "Analysis: OK"
            if result['is_outlier']:
                status_line = "STATUS: OUTLIER (Excluded)"

            lines = [
                f"Object ID: {result['id']}",
                f"Coords(x,y,z): ({result['cx']}, {result['cy']}, {result['cz']})",
                f"Volume(voxels): {result['area']}",
                f"Trend: {result['trend_class']}",
                f"Size: {result['size_class']}",
                status_line
            ]
            
            for line in lines:
                p = tf.add_paragraph()
                p.text = line
                p.font.size = Inches(0.25)
                
                if "Trend" in line and result['trend_class'] != "No trend":
                    p.font.bold = True
                
                if "OUTLIER" in line:
                    p.font.color.rgb = RGBColor(255, 0, 0)
                    p.font.bold = True

    # --- SUMMARY SLIDES ---
    print("Generating Summary Slides...")
    
    trend_types = ["Colocalized", "Around", "Anticolocalized", "No trend"]
    size_types = ["small", "large"] 
    
    for trend in trend_types:
        for size in size_types:
            key = (trend, size)
            
            if key in groups and len(groups[key]['r']) > 0:
                slide = prs.slides.add_slide(blank_layout)
                
                avg_r = np.mean(groups[key]['r'], axis=0)
                avg_g = np.mean(groups[key]['g'], axis=0)
                
                plot_stream = create_plot_image(avg_r, avg_g, np.arange(11), "AVERAGE")
                slide.shapes.add_picture(plot_stream, Inches(0.5), Inches(1), height=Inches(5))
                
                tx = slide.shapes.add_textbox(Inches(8), Inches(2.5), Inches(5), Inches(2))
                tf = tx.text_frame
                tf.text = f"Group Average\nTrend: {trend}\nSize: {size}\nCount (n): {len(groups[key]['r'])}"
                for p in tf.paragraphs:
                    p.font.size = Inches(0.25)
                    p.font.bold = True

    # --- SAVE OUTPUT ---
    while True:
        try:
            prs.save(output_pptx)
            print(f"Saved successfully: {output_pptx}")
            break
        except PermissionError:
            input(f"File {output_pptx} is OPEN. Close it and press ENTER...")


if __name__ == "__main__":
    # --- CONFIGURATION ---
    # Update these paths to point to your 3D TIFF files
    
    base_path = r"/media/kyoung/Elements/20260101 PFKL-mCherry_EGFP-GAPDH - High Salt Conc - 24.1 degree - LLSM/Deconv/Cell2 - 180mM/3_NoGlucose60min"
    
    if os.path.exists(base_path):
        main(
            os.path.join(base_path, "Cell2_ex560em593_250mW_200ms_NoGlucose60min_d62_corrected.tif"), 
            os.path.join(base_path, "Cell2_ex488em514_230mW_200ms_NoGlucose60min_d16.tif"), 
            os.path.join(base_path, "Cell2_ex560em593_250mW_200ms_NoGlucose60min_d62_corrected-1.segmentation.tif"),     
            os.path.join(base_path, "Cell2_NoGlucose60min_3D_Report.pptx")
        )
    else:
        print("Please configure the 'base_path' in the script.")