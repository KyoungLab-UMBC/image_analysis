import matplotlib
# Force Agg backend to prevent GUI errors
matplotlib.use('Agg') 

# Plotting imports
from matplotlib.figure import Figure
from matplotlib.backends.backend_agg import FigureCanvasAgg

import numpy as np
import tifffile
from skimage import measure, draw
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
import io
import re
import warnings
import os
from pathlib import Path
from read_roi import read_roi_zip
import concurrent.futures
from PIL import Image 
import pandas as pd

# --- CUSTOM MODULE IMPORTS ---
import util.background_correction as bg_tools
import util.trend_analysis as trend_tools
import distance_analysis as dis_tool
import util.config as cfg
import util.roi_to_mask as roi_tools
import util.region_separation as rs

# Suppress runtime warnings
warnings.filterwarnings("ignore", category=RuntimeWarning)
warnings.filterwarnings("ignore", category=UserWarning)

# --- HELPER FUNCTIONS ---

def create_plot_image(red_data, green_data, x_axis, obj_id, r_name, g_name):
    fig = Figure(figsize=(6, 4))
    canvas = FigureCanvasAgg(fig)
    
    ax1 = fig.add_subplot(111)
    
    ax1.set_xlabel('Radius(pixel)')
    ax1.set_xlim(0, 20)
    ax1.set_xticks(range(0, 21, 5))
    
    # Left Y: Red
    ax1.set_ylabel(f'Norm. {r_name}', color='red')
    ax1.plot(x_axis, red_data, color='red', linewidth=1.5)
    ax1.tick_params(axis='y', labelcolor='red')

    # Right Y: Green
    ax2 = ax1.twinx() 
    ax2.set_ylabel(f'Norm. {g_name}', color='green')
    ax2.plot(x_axis, green_data, color='green', linewidth=1.5)
    ax2.tick_params(axis='y', labelcolor='green')

    if obj_id != "AVERAGE":
        ax1.set_title(f"Object: {obj_id}")
    
    fig.tight_layout()
    
    img_stream = io.BytesIO()
    canvas.print_png(img_stream)
    img_stream.seek(0)
    return img_stream

def create_composite_thumbnail(norm_r, norm_g, norm_m=None):
    h, w = norm_r.shape
    rgb = np.zeros((h, w, 3), dtype=np.uint8)
    
    def scale_to_uint8(arr):
        vmin = arr.min()
        vmax = arr.max()
        if vmax - vmin < 0.00001:
            return np.zeros_like(arr, dtype=np.uint8)
        scaled = (arr - vmin) / (vmax - vmin) * 255
        return scaled.astype(np.uint8)

    if norm_m is not None:
        # MITO MODE: R=Mito, G=PFKL(RedData), B=GAPDH(GreenData)
        rgb[..., 0] = scale_to_uint8(norm_m) 
        rgb[..., 1] = scale_to_uint8(norm_r) 
        rgb[..., 2] = scale_to_uint8(norm_g) 
    else:
        # STANDARD MODE: R=PFKL, G=GAPDH, B=Blank
        rgb[..., 0] = scale_to_uint8(norm_r) 
        rgb[..., 1] = scale_to_uint8(norm_g) 
        
    img = Image.fromarray(rgb)
    img_stream = io.BytesIO()
    img.save(img_stream, format='PNG')
    img_stream.seek(0)
    return img_stream

def process_single_object(data_package):
    # --- CHANGED: Unpack new r_name and g_name arguments ---
    (obj_id, cy, cx, area, crop_r, crop_g, crop_m, crop_mask, mito_mode, dist_val, r_name, g_name, generate_slides) = data_package

    h, w = crop_r.shape
    center_y, center_x = h // 2, w // 2

    # Create binary mask for this specific object
    obj_mask = (crop_mask == obj_id)
    
    # --- CHANGED: Calculate Max Intensity inside for IO PC ---
    r_in = np.max(crop_r[obj_mask])
    g_in = np.max(crop_g[obj_mask])

    # Sum raw intensities within the object
    raw_int_r = np.sum(crop_r[obj_mask])
    raw_int_g = np.sum(crop_g[obj_mask])

    # --- 1. Background Estimation (Red/Green only) ---
    bg_r = bg_tools.estimate_background_rolling_ball(crop_r, radius=10)
    bg_r[bg_r == 0] = 1
    norm_img_r = crop_r.astype(float) / bg_r
    
    bg_g = bg_tools.estimate_background_rolling_ball(crop_g, radius=10)
    bg_g[bg_g == 0] = 1
    norm_img_g = crop_g.astype(float) / bg_g
    
    # Mito is now a mask, just normalize for display (0-255)
    norm_img_m = None
    if mito_mode and crop_m is not None:
        norm_img_m = crop_m.astype(float) 
    
    # --- 2. Radial Profile ---
    y, x = np.indices((h, w))
    radii_map = np.sqrt((x - center_x)**2 + (y - center_y)**2)
    
    profile_r = []
    profile_g = []
    
    for r in range(21):
        ring_mask = (radii_map >= r) & (radii_map < r + 1)
        valid_mask = (crop_mask == 0) | (crop_mask == obj_id)
        final_mask = ring_mask & valid_mask
        
        vals_r = norm_img_r[final_mask]
        vals_g = norm_img_g[final_mask]
        
        profile_r.append(np.mean(vals_r) if len(vals_r) > 0 else 0)
        profile_g.append(np.mean(vals_g) if len(vals_g) > 0 else 0)

    final_r = np.array(profile_r)
    final_g = np.array(profile_g)
    
    # --- 3. Classification ---
    size_class = "small glucosome" if area <= 14 else "large glucosome"
    trend_class, r_value, sigma, pc_r, pc_g = trend_tools.analyze_profile_trend(final_r, final_g)

    # --- CHANGED: Outlier Detection ---
    is_outlier = False
    # Check Max Intensity > 30 OR Peak not at center (index 0 or 1)
    if np.max(final_r) > 30 or np.max(final_g) > 30 or np.argmax(final_r) not in [0, 1]:
        is_outlier = True
    
    # --- CHANGED: Distance Based Classification ---
    mito_status = "N/A"
    if mito_mode:
        # Check distance threshold (275 nm)
        if dist_val is not None and dist_val <= 275:
            mito_status = "Attached"
        else:
            mito_status = "Not Attached"
    
    # --- 4. Plotting ---
    plot_bytes = b""
    thumb_bytes = b""
    
    if generate_slides:
        x_vals = np.arange(0, 21)
        # --- CHANGED: Pass names and mode to plotter ---
        plot_bytes = create_plot_image(final_r, final_g, x_vals, obj_id, r_name, g_name).getvalue()
        # Thumbnails: Red=Mito, Green=PFKL, Blue=GAPDH
        thumb_bytes = create_composite_thumbnail(norm_img_r, norm_img_g, norm_img_m).getvalue()
    
    return {
        "id": obj_id,
        "cy": cy, "cx": cx,
        "area": area,
        "norm_r": final_r,
        "norm_g": final_g,
        "size_class": size_class,
        "trend_class": trend_class,
        "r_value": r_value,
        "sigma": sigma,
        "pc_r": pc_r,   
        "pc_g": pc_g,   
        "r_in": r_in,           
        "g_in": g_in,           
        "raw_int_r": raw_int_r, 
        "raw_int_g": raw_int_g,
        "dist_val": dist_val, 
        "mito_status": mito_status, 
        "is_outlier": is_outlier,
        "plot_bytes": plot_bytes,
        "thumb_bytes": thumb_bytes
    }

# --- MAIN CONTROLLER ---

def processing_cell_roi(user_path, img_r, img_g, mito_mode, r_name, g_name, roi_name, roi_data, generate_slides):
# Load Mito Mask if needed
    img_m = None
    dist_map = {}
    mito_path = None
    if mito_mode:
        mito_path = user_path / f"{cfg.MITO_RAW.stem}_{roi_name}_mask.tif"
    # Construct paths using stem and roi_name (cell_num)
    r_mask_path = user_path / f"{cfg.R_RAW.stem}_{roi_name}_mask.tif"
    r_mask = tifffile.imread(r_mask_path)
    output_pptx = user_path / f"{cfg.R_RAW.stem}_{roi_name}.pptx"
    print(f"Analyzing ROI: {roi_name} | Output: {output_pptx.name}")
    img_r = bg_tools.estimate_background_rolling_ball(img_r, radius=10, create_background=False, use_paraboloid=True)
    img_g = bg_tools.estimate_background_rolling_ball(img_g, radius=10, create_background=False, use_paraboloid=True)

    if mito_mode and mito_path:
        print("Loading Mito Mask and calculating distances...")
        img_m = tifffile.imread(mito_path) # Load as mask
    
        # --- DISTANCE ANALYSIS START ---
        # 1. Get objects from PFKL Mask (Source)
        # We explicitly pass is_2d=True
        pfkl_objs, _ = dis_tool.get_object_data(r_mask_path, True)
    
        # 2. Get objects from Mito Mask (Target Cloud)
        _, mito_cloud = dis_tool.get_object_data(mito_path, True)
    
        # 3. Calculate Distances
        dists = dis_tool.analyze_distances(pfkl_objs, mito_cloud, True, s_xy=107, s_z=107)
        
        # 4. Map Label -> Distance
        # dis_tool returns distances in the same order as pfkl_objs
        for obj, d in zip(pfkl_objs, dists):
            dist_map[obj['label']] = d
        # --- DISTANCE ANALYSIS END ---
    # --- NEW: Calculate Total Intensity inside the ROI (Denominator) ---
    # Use the imported roi_to_mask function
    roi_mask_full = roi_tools.roi_to_mask(roi_data, img_r.shape)
    
    # Calculate Total Intensity of R and G inside the ROI
    total_roi_int_r = np.sum(img_r[roi_mask_full])
    total_roi_int_g = np.sum(img_g[roi_mask_full])
    
    if total_roi_int_r == 0: total_roi_int_r = 1 # Avoid div/0
    if total_roi_int_g == 0: total_roi_int_g = 1

    # --- CHANGED: Calculate background and delegate to region_separation ---
    print("Calculating background and separating regions...")
    
    # Delegate logic to region_separation.py
    r_bg, g_bg, mask_hh, mask_lh, mask_hl, mask_ll = rs.separate_regions_and_bg(img_r, img_g, roi_mask_full)
    
    # Save the 4 quadrant masks to a zip file
    regions_zip_path = user_path / f"{cfg.R_RAW.stem}_{roi_name}_regions.zip"
    rs.save_masks_to_tifs(regions_zip_path, mask_hh, mask_lh, mask_hl, mask_ll)

    def get_out_val(bg_img, mask):
        vals = bg_img[mask].flatten()
        # Sample rate 1/100
        sampled = vals[::100]
        if len(sampled) == 0: return 1.0       
        # Sort values in ascending order (lowest to highest)
        sampled = np.sort(sampled)        
        # Mean of first 10% (lowest values in the sampled group)
        idx = max(1, int(len(sampled) * 0.1))
        return np.mean(sampled[:idx])

    # 2. Calculate the 'out' values for each of the 4 masks
    hh_r_out = get_out_val(r_bg, mask_hh)
    hh_g_out = get_out_val(g_bg, mask_hh)
    
    lh_r_out = get_out_val(r_bg, mask_lh)
    lh_g_out = get_out_val(g_bg, mask_lh)
    
    hl_r_out = get_out_val(r_bg, mask_hl)
    hl_g_out = get_out_val(g_bg, mask_hl)
    
    ll_r_out = get_out_val(r_bg, mask_ll)
    ll_g_out = get_out_val(g_bg, mask_ll)

    print(f"High {r_name} & High {g_name} (hh): {hh_r_out:.2f}, {hh_g_out:.2f}")
    print(f"Low {r_name} & High {g_name} (lh): {lh_r_out:.2f}, {lh_g_out:.2f}")
    print(f"High {r_name} & Low {g_name} (hl): {hl_r_out:.2f}, {hl_g_out:.2f}")
    print(f"Low {r_name} & Low {g_name} (ll): {ll_r_out:.2f}, {ll_g_out:.2f}")
    # ------------------------------------------------------------------

    if len(r_mask.shape) > 2: r_mask = r_mask[..., 0]
    
    # ... (Mask Labeling logic remains the same) ...
    # Re-running label logic just to be safe for regionprops
    unique_vals = np.unique(r_mask)
    if len(unique_vals) <= 2 or r_mask.max() == 255:
        labeled_mask = measure.label(r_mask > 0)
    else:
        labeled_mask = r_mask

    props = measure.regionprops(labeled_mask)
    print(f"Found {len(props)} objects. Starting analysis...")

    tasks = []
    h_img, w_img = img_r.shape
    
    for prop in props:
        cy, cx = map(int, prop.centroid)
        r1, r2 = max(0, cy - 30), min(h_img, cy + 30 + 1)
        c1, c2 = max(0, cx - 30), min(w_img, cx + 30 + 1)
        
        # Crop Mito Mask if available
        crop_m = None
        if mito_mode and img_m is not None:
            crop_m = img_m[r1:r2, c1:c2]

        # Retrieve pre-calculated distance for this object
        # Default to a large number (99999) if not found, so it classifies as "Not Attached"
        dist_val = dist_map.get(prop.label, 99999) if mito_mode else None

        tasks.append((
            prop.label, cy, cx, prop.area,
            img_r[r1:r2, c1:c2],
            img_g[r1:r2, c1:c2],
            crop_m,                   
            labeled_mask[r1:r2, c1:c2],
            mito_mode,
            dist_val,
            r_name, 
            g_name,
            generate_slides  
        ))

    results = []
    
    # ThreadPool
    max_workers = min(32, os.cpu_count())
    print(f"Processing with ThreadPool ({max_workers} workers)...")
    
    with concurrent.futures.ThreadPoolExecutor(max_workers=max_workers) as executor:
        for i, res in enumerate(executor.map(process_single_object, tasks)):
            results.append(res)
            if (i + 1) % 50 == 0:
                print(f"Analyzed {i+1}/{len(tasks)}")
    
    # --- CHANGED: Process IO PC for Colocalized Objects ---
    for res in results:
        if res['trend_class'] == "Colocalized":
            # Ensure coordinates are ints for array indexing
            cy, cx = int(res['cy']), int(res['cx'])
            
            # Ensure coordinates are within background image bounds
            cy_c = min(max(cy, 0), r_bg.shape[0]-1)
            cx_c = min(max(cx, 0), r_bg.shape[1]-1)
            
            # Use center position to decide which of the 4 masks it belongs to
            if mask_hh[cy_c, cx_c]:
                r_out, g_out = hh_r_out, hh_g_out
            elif mask_lh[cy_c, cx_c]:
                r_out, g_out = lh_r_out, lh_g_out
            elif mask_hl[cy_c, cx_c]:
                r_out, g_out = hl_r_out, hl_g_out
            else: # mask_ll
                r_out, g_out = ll_r_out, ll_g_out
                
            res['r_out'] = r_out
            res['g_out'] = g_out
            
            res['pc_r_io'] = res['r_in'] / r_out if r_out != 0 else 0
            res['pc_g_io'] = res['g_in'] / g_out if g_out != 0 else 0
            res['pc_ratio_io'] = res['pc_g_io'] / res['pc_r_io'] if res['pc_r_io'] != 0 else 0

    # --- PPT GENERATION ---
    groups = {}
    if generate_slides:
        print("Generating PowerPoint Report...")
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank_layout = prs.slide_layouts[6]

    for res in results:
        # --- CHANGED: Include mito_status in the key if in mito_mode ---
        if mito_mode:
            k = (res['trend_class'], res['size_class'], res['mito_status'])
        else:
            k = (res['trend_class'], res['size_class'])
        
        # --- CHANGED: Only add to summary groups if NOT outlier ---
        if not res['is_outlier']:
            if k not in groups: groups[k] = {'r': [], 'g': []}
            groups[k]['r'].append(res['norm_r'])
            groups[k]['g'].append(res['norm_g'])
        
        if generate_slides:
            slide = prs.slides.add_slide(blank_layout)
            
            graph_stream = io.BytesIO(res['plot_bytes'])
            slide.shapes.add_picture(graph_stream, left=Inches(0.5), top=Inches(1), width=Inches(7.5), height=Inches(5))
            
            thumb_stream = io.BytesIO(res['thumb_bytes'])
            slide.shapes.add_picture(thumb_stream, left=Inches(8), top=Inches(0.5), width=Inches(3), height=Inches(3))
            
            tx = slide.shapes.add_textbox(left=Inches(8), top=Inches(3.5), width=Inches(5), height=Inches(3))
            tf = tx.text_frame
            tf.word_wrap = True

            # --- CHANGED: Add Colored Legend Lines ---
            if mito_mode:
                # 1. Mitochondria (Red)
                p = tf.add_paragraph()
                p.text = "Mitochondria"
                p.font.color.rgb = RGBColor(255, 0, 0)
                p.font.size = Inches(0.25)
                p.font.bold = True
                
                # 2. r_name (Green)
                p = tf.add_paragraph()
                p.text = r_name
                p.font.color.rgb = RGBColor(0, 255, 0)
                p.font.size = Inches(0.25)
                p.font.bold = True

                # 3. g_name (Blue)
                p = tf.add_paragraph()
                p.text = g_name
                p.font.color.rgb = RGBColor(0, 0, 255)
                p.font.size = Inches(0.25)
                p.font.bold = True
            else:
                # 1. r_name (Red)
                p = tf.add_paragraph()
                p.text = r_name
                p.font.color.rgb = RGBColor(255, 0, 0)
                p.font.size = Inches(0.25)
                p.font.bold = True

                # 2. g_name (Green)
                p = tf.add_paragraph()
                p.text = g_name
                p.font.color.rgb = RGBColor(0, 255, 0)
                p.font.size = Inches(0.25)
                p.font.bold = True
            
            # --- Standard Lines ---
            lines = [
                f"Object ID: {res['id']}",
                f"Coords(x,y): ({res['cx']}, {res['cy']})",
                f"Size Class: {res['size_class']} ({res['area']} px)",
                f"Sigma: {res['sigma']:.2f}",
                f"Trend: {res['trend_class']}",
                f"Pearson coefficient: {res['r_value']:.2f}"
            ]
            # Add Mito Status line if applicable
            if mito_mode:
                lines.append(f"Mito: {res['mito_status']}")

            # --- CHANGED: Add Outlier Warning Label ---
            if res['is_outlier']:
                lines.append("!!! OUTLIER - EXCLUDED !!!")

            for line in lines:
                p = tf.add_paragraph()
                p.text = line
                p.font.size = Inches(0.25)
                if "Trend" in line and res['trend_class'] != "No trend":
                    p.font.bold = True
                # Highlight Attached status
                if "Mito" in line and "Attached" in line and "Not" not in line:
                    p.font.bold = True

    # =========================================================================
    # CONSOLIDATED CATEGORY & GROUP LOGIC
    # =========================================================================
    
    # 1. Define Standard Types
    trend_types = ["Colocalized", "Around", "Anticolocalized", "No trend"]
    size_types = ["small glucosome", "large glucosome"]
    mito_types = ["Attached", "Not Attached"] if mito_mode else [None]

    # 2. Build Categories List (Used for both Summary Slides and Excel Sheets)
    categories = []
    if mito_mode:
        for trend in trend_types:
            for size in size_types:
                for mito in mito_types:
                    categories.append((trend, size, mito))
    else:
        for trend in trend_types:
            for size in size_types:
                categories.append((trend, size))

    # =========================================================================
    # 1. GENERATE SUMMARY SLIDES
    # =========================================================================
    if generate_slides:
        print("Generating Summary Slides...")
        
        for cat_tuple in categories:
            # Construct Key based on mode
            if mito_mode:
                trend, size, mito = cat_tuple
                key = (trend, size, mito)
                title_str = f"Group Average\n{trend}\n{size}\n{mito}"
            else:
                trend, size = cat_tuple
                key = (trend, size)
                title_str = f"Group Average\n{trend}\n{size}"

            # CHECK IF KEY EXISTS AND HAS DATA (Non-outliers are already filtered in 'groups' dict)
            if key in groups and len(groups[key]['r']) > 0:
                count = len(groups[key]['r'])
                title_str += f"\n(n={count})"

                slide = prs.slides.add_slide(blank_layout)
                avg_r = np.mean(groups[key]['r'], axis=0)
                avg_g = np.mean(groups[key]['g'], axis=0)
                
                plot_stream = create_plot_image(avg_r, avg_g, np.arange(21), "AVERAGE", r_name, g_name)
                slide.shapes.add_picture(plot_stream, Inches(0.5), Inches(1), height=Inches(5))
                
                tx = slide.shapes.add_textbox(Inches(8), Inches(2.5), Inches(5), Inches(2))
                tx.text_frame.text = title_str
                for p in tx.text_frame.paragraphs:
                    p.font.bold = True
                    p.font.size = Inches(0.25)

        prs.save(output_pptx)
        print(f"Saved PowerPoint: {output_pptx}")

    # =========================================================================
    # 2. GENERATE EXCEL REPORT
    # =========================================================================
    print("Generating Excel Report...")
    output_xlsx = output_pptx.with_suffix(".xlsx")
    
    # --- Calculate ROI Area and Means ---
    roi_area = np.sum(roi_mask_full)
    if roi_area == 0: roi_area = 1
    mean_roi_int_r = total_roi_int_r / roi_area
    mean_roi_int_g = total_roi_int_g / roi_area
    glucosome_num = len(results)

    # --- Offset and Time Point Logic ---
    folder_name = user_path.name
    start_col_idx = 0
    prefix_num = 0
    
    match_prefix = re.match(r'^(\d+)_', folder_name)
    
    if folder_name == "0" or "Cell" in folder_name:
        start_col_idx = 0
        prefix_num = 0
    elif match_prefix:
        prefix_num = int(match_prefix.group(1))
        start_col_idx = prefix_num * 7 

    if start_col_idx == 0:
        time_point = prefix_num
    else:
        time_point = folder_name

    try:
        with pd.ExcelWriter(output_xlsx, engine='openpyxl') as writer:
            
            # --- Construct Custom Image ID ---
            path_str = str(user_path) 
            date_part = "00000000"
            date_match = re.search(r'(\d{8})', path_str)
            if date_match:
                date_part = date_match.group(1)
            
            cell_part = "0"
            cell_match = re.search(r'[Cc]ell(\d+)', path_str)
            if cell_match:
                cell_part = cell_match.group(1)
                
            image_id = f"{date_part}-{cell_part}-{roi_name}-{time_point}"

            # --- 1. Prepare Summary Dictionaries ---
            summary_main = {
                "Image Name": image_id,
                "Area/px": roi_area,
                "Glucosome num": glucosome_num,
                f"Mean {r_name} int": mean_roi_int_r,
                f"Total {r_name} int": total_roi_int_r,
                f"Mean {g_name} int": mean_roi_int_g,
                f"Total {g_name} int": total_roi_int_g
            }

            summary_num = {}
            summary_ratio_count = {}
            summary_ratio_area = {}
            summary_ratio_r = {}
            summary_ratio_g = {}
            
            summary_num["Image Name"] = image_id
            summary_ratio_count["Image Name"] = image_id
            summary_ratio_area["Image Name"] = image_id
            summary_ratio_r["Image Name"] = image_id
            summary_ratio_g["Image Name"] = image_id
            
            # Calculate Total Valid Objects (for Count Ratio)
            valid_objects = [r for r in results if not r['is_outlier']]
            total_valid_count = len(valid_objects) if len(valid_objects) > 0 else 1

            for cat_tuple in categories:
                if mito_mode:
                    trend, size, mito = cat_tuple
                    group_items = [r for r in valid_objects if r['trend_class'] == trend and r['size_class'] == size and r['mito_status'] == mito]
                    group_name = f"{trend} {size} {mito}"
                else:
                    trend, size = cat_tuple
                    group_items = [r for r in valid_objects if r['trend_class'] == trend and r['size_class'] == size]
                    group_name = f"{trend} {size}"
                
                count = len(group_items)
                summary_num[group_name] = count
                summary_ratio_count[group_name] = count / total_valid_count

                group_sum_area = sum(item['area'] for item in group_items)
                summary_ratio_area[group_name] = group_sum_area / roi_area
                
                group_sum_r = sum(item['raw_int_r'] for item in group_items)
                summary_ratio_r[group_name] = group_sum_r / total_roi_int_r
                
                group_sum_g = sum(item['raw_int_g'] for item in group_items)
                summary_ratio_g[group_name] = group_sum_g / total_roi_int_g

            # --- 2. Write Summary Sheets ---
            # Main Summary (startcol = prefix_num * 8)
            pd.DataFrame([summary_main]).to_excel(writer, sheet_name="Summary", index=False, startcol=prefix_num * 8)

            # Other Summaries (startcol = prefix_num * 18)
            if mito_mode:
                pd.DataFrame([summary_num]).to_excel(writer, sheet_name="Summary_num", index=False, startcol=prefix_num * 18)
                pd.DataFrame([summary_ratio_count]).to_excel(writer, sheet_name="Summary_ratio", index=False, startcol=prefix_num * 18)
                pd.DataFrame([summary_ratio_area]).to_excel(writer, sheet_name="Summary_area_ratio", index=False, startcol=prefix_num * 18)
                pd.DataFrame([summary_ratio_r]).to_excel(writer, sheet_name=f"Summary_{r_name}_ratio", index=False, startcol=prefix_num * 18)
                pd.DataFrame([summary_ratio_g]).to_excel(writer, sheet_name=f"Summary_{g_name}_ratio", index=False, startcol=prefix_num * 18)
            else:
                pd.DataFrame([summary_num]).to_excel(writer, sheet_name="Summary_num", index=False, startcol=prefix_num * 10)
                pd.DataFrame([summary_ratio_count]).to_excel(writer, sheet_name="Summary_ratio", index=False, startcol=prefix_num * 10)
                pd.DataFrame([summary_ratio_area]).to_excel(writer, sheet_name="Summary_area_ratio", index=False, startcol=prefix_num * 10)
                pd.DataFrame([summary_ratio_r]).to_excel(writer, sheet_name=f"Summary_{r_name}_ratio", index=False, startcol=prefix_num * 10)
                pd.DataFrame([summary_ratio_g]).to_excel(writer, sheet_name=f"Summary_{g_name}_ratio", index=False, startcol=prefix_num * 10)

            # --- 3. Generate Data Sheets ---
            for cat_tuple in categories:
                if mito_mode:
                    trend, size, mito = cat_tuple
                    subset = [r for r in results if r['trend_class'] == trend and r['size_class'] == size and r['mito_status'] == mito and not r['is_outlier']]
                    t_abbr = trend[:5] 
                    s_abbr = "Sm" if "small" in size else "Lg"
                    m_abbr = "Att" if mito == "Attached" else "NotAtt"
                    sheet_name = f"{t_abbr} {s_abbr} {m_abbr}"
                else:
                    trend, size = cat_tuple
                    subset = [r for r in results if r['trend_class'] == trend and r['size_class'] == size and not r['is_outlier']]
                    sheet_name = f"{trend} {size}"[:31]
                
                if not subset: continue

                # Subset standard sheets (startcol = prefix_num * 46)
                header_row = [f"Normalized {g_name} {time_point}"] + list(range(21)) + ["", f"Normalized {r_name} {time_point}"] + list(range(21))
                data_rows = [header_row]
                
                for item in subset:
                    row_data = [item['id']] + list(item['norm_g']) + ["", item['id']] + list(item['norm_r'])
                    data_rows.append(row_data)
                
                df_profile = pd.DataFrame(data_rows)
                df_profile.to_excel(writer, sheet_name=sheet_name, index=False, header=False, startcol=prefix_num * 46)

                # --- New Logic for "Colocalized" Partition Coefficient Sheets ---
                if trend == "Colocalized":
                    pc_sheet_name = f"{sheet_name}_pc"[:31]
                    
                    # Columns updated (startcol = prefix_num * 15)
                    pc_headers = [
                        f"Obj ID {time_point}", "Sigma", "Size (px)", 
                        f"PC {r_name} Radius", f"PC {g_name} Radius", "PC Radius Ratio", 
                        f"{r_name} in", f"{r_name} out", f"{r_name} PC IO",
                        f"{g_name} in", f"{g_name} out", f"{g_name} PC IO",
                        "PC IO ratio (G/R)", "Mito Dist"
                    ]
                    pc_rows = []
                    
                    for item in subset:
                        ratio = 0
                        if item['pc_r'] is not None and item['pc_r'] != 0 and item['pc_g'] is not None:
                            ratio = item['pc_g'] / item['pc_r']
                        
                        d_val = item['dist_val'] if mito_mode else ""
                        
                        row = [
                            item['id'],
                            item['sigma'],
                            item['area'],
                            item['pc_r'],
                            item['pc_g'],
                            ratio,
                            item.get('r_in', ''),
                            item.get('r_out', ''),
                            item.get('pc_r_io', ''),
                            item.get('g_in', ''),
                            item.get('g_out', ''),
                            item.get('pc_g_io', ''),
                            item.get('pc_ratio_io', ''),
                            d_val
                        ]
                        pc_rows.append(row)
                    
                    df_pc = pd.DataFrame(pc_rows, columns=pc_headers)
                    df_pc.to_excel(writer, sheet_name=pc_sheet_name, index=False, startcol=prefix_num * 15)

        print(f"Saved Excel: {output_xlsx}")
        
    except Exception as e:
        print(f"Failed to save Excel: {e}")
        import traceback
        traceback.print_exc()


def main():    
    if not cfg.R_PATH.exists():
        print(f"Path not found: {cfg.R_PATH}")
    else:
        # Initialize ImageJ once
        bg_tools.init_imagej()
        print(f"Running analysis for: {cfg.R_NAME} vs {cfg.G_NAME}")
        user_path = cfg.USER_PATH
        mito_mode = cfg.MITO_ANALYSIS
        print("Loading images...")
        try:
            img_r = tifffile.imread(cfg.R_PATH)
            img_g = tifffile.imread(cfg.G_PATH)      
            rois = roi_tools.load_roi_file(cfg.R_RAW)
            if not rois:
                print(f"Error: rois.zip not found in {cfg.USER_PATH}")
            else:             
                for roi_name, roi_data in rois.items():
                    processing_cell_roi(
                        user_path,
                        img_r, 
                        img_g,
                        mito_mode, 
                        r_name=cfg.R_NAME, 
                        g_name=cfg.G_NAME,
                        roi_name=roi_name,
                        roi_data=roi_data,
                        generate_slides = cfg.GENERATE_SLIDES
                    )
        except Exception as e:
            print(f"Error loading or analyzing distances: {e}")
            import traceback
            traceback.print_exc()
            return

if __name__ == "__main__":
    main()